"""
文档解析工具模块
---------------
支持解析 PDF、Word(.docx)、TXT 格式的运维文档，
提取纯文本内容，并进行清洗处理。

清洗步骤：
1. 移除多余空白（连续换行、空格）
2. 统一编码为 UTF-8
3. 去除不可见控制字符
"""

import os
import re


def parse_pdf(file_path: str) -> str:
    """
    解析PDF文件，提取纯文本内容。

    Args:
        file_path: PDF文件的绝对路径

    Returns:
        提取到的文本字符串
    """
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        text_parts = []
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
        return "\n".join(text_parts)
    except ImportError:
        raise ImportError("请安装 PyPDF2 库: pip install PyPDF2")
    except Exception as e:
        raise RuntimeError(f"PDF解析失败 [{file_path}]: {e}")


def parse_docx(file_path: str) -> str:
    """
    解析Word(.docx)文件，提取纯文本内容。

    Args:
        file_path: docx文件的绝对路径

    Returns:
        提取到的文本字符串
    """
    try:
        from docx import Document
        doc = Document(file_path)
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        return "\n".join(text_parts)
    except ImportError:
        raise ImportError("请安装 python-docx 库: pip install python-docx")
    except Exception as e:
        raise RuntimeError(f"Word解析失败 [{file_path}]: {e}")


def parse_txt(file_path: str) -> str:
    """
    读取TXT文件内容，自动尝试常见编码。

    Args:
        file_path: TXT文件的绝对路径

    Returns:
        文件文本内容
    """
    # 按优先级尝试不同编码
    for encoding in ["utf-8", "gbk", "gb2312", "latin-1"]:
        try:
            with open(file_path, "r", encoding=encoding) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    # 最终兜底：忽略无法解码的字符
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def parse_pptx(file_path: str) -> str:
    """
    解析 PowerPoint(.pptx) 文件，提取所有幻灯片文本和表格。

    Args:
        file_path: pptx 文件的绝对路径

    Returns:
        提取到的文本字符串
    """
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        txt = para.text.strip()
                        if txt:
                            slide_text.append(txt)
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(" | ".join(cells))
                    if rows:
                        slide_text.append("\n".join(rows))
            if slide_text:
                parts.append(f"[Slide {i}]\n" + "\n".join(slide_text))
        return "\n\n".join(parts)
    except ImportError:
        raise ImportError("请安装 python-pptx 库: pip install python-pptx")
    except Exception as e:
        raise RuntimeError(f"PowerPoint 解析失败 [{file_path}]: {e}")


def parse_html(file_path: str) -> str:
    """
    解析 HTML 文件，提取文本内容（去除脚本和样式标签）。

    Args:
        file_path: HTML 文件的绝对路径

    Returns:
        提取到的文本字符串
    """
    try:
        from bs4 import BeautifulSoup
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            soup = BeautifulSoup(f.read(), "lxml")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)
    except ImportError:
        raise ImportError("请安装 beautifulsoup4 + lxml: pip install beautifulsoup4 lxml")
    except Exception as e:
        raise RuntimeError(f"HTML 解析失败 [{file_path}]: {e}")


def parse_document(file_path: str, use_mineru: bool = None) -> str:
    """
    通用文档解析入口：根据文件后缀自动选择解析器。

    支持格式：.pdf / .docx / .pptx / .html / .htm / .txt / 图片格式

    Args:
        file_path: 文档文件路径
        use_mineru: 是否尝试 MinerU 解析。None=按 config 配置决定。

    Returns:
        清洗后的纯文本内容
    """
    # MinerU 路径
    if use_mineru is None:
        from utils.config import config
        use_mineru = config.get("knowledge_base.mineru.enabled", False)

    if use_mineru:
        try:
            from utils.doc_parser_v2 import parse_document_v2
            return parse_document_v2(file_path, use_mineru=True)
        except Exception:
            pass  # 回退到传统解析

    ext = os.path.splitext(file_path)[1].lower()

    from utils.ocr import SUPPORTED_IMAGE_EXTS

    if ext == ".pdf":
        raw_text = parse_pdf(file_path)
    elif ext == ".docx":
        raw_text = parse_docx(file_path)
    elif ext == ".pptx":
        raw_text = parse_pptx(file_path)
    elif ext in (".htm", ".html"):
        raw_text = parse_html(file_path)
    elif ext == ".txt":
        raw_text = parse_txt(file_path)
    elif ext in SUPPORTED_IMAGE_EXTS:
        raw_text = parse_image(file_path)
    else:
        raise ValueError(
            f"不支持的文件格式: {ext}"
            f"（支持 .pdf / .docx / .pptx / .html / .txt / 图片格式）"
        )

    return clean_text(raw_text)


def parse_image(file_path: str) -> str:
    """
    使用 OCR 从图片中提取文字。
    支持 jpg/png/bmp/tiff/webp 等常见格式。

    Args:
        file_path: 图片文件路径

    Returns:
        OCR 识别出的文字内容
    """
    from utils.ocr import extract_text, SUPPORTED_IMAGE_EXTS

    ext = os.path.splitext(file_path)[1].lower()
    if ext not in SUPPORTED_IMAGE_EXTS:
        raise ValueError(f"不支持的图片格式: {ext}（支持 {', '.join(sorted(SUPPORTED_IMAGE_EXTS))}）")

    return extract_text(file_path)


def clean_text(text: str) -> str:
    """
    清洗文本：去除多余空白、控制字符，统一格式。

    Args:
        text: 原始文本

    Returns:
        清洗后的文本
    """
    if not text:
        return ""

    # 去除不可见的控制字符（保留换行和制表符）
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]", "", text)

    # 将连续3个以上的换行压缩为2个
    text = re.sub(r"\n{3,}", "\n\n", text)

    # 将连续空白行压缩
    text = re.sub(r"[ \t]+", " ", text)

    # 去除首尾空白
    text = text.strip()

    return text


def split_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list:
    """
    将长文本按固定大小切分为块，块之间有重叠以保持上下文连贯。
    中文优化：按句号、换行等自然边界优先切分。

    Args:
        text: 待切分的文本
        chunk_size: 每块的字数上限
        overlap: 相邻块之间的重叠字数

    Returns:
        文本块列表
    """
    if len(text) <= chunk_size:
        return [text]

    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]

        # 尝试在句号处断开，使语义更完整
        if end < len(text):
            # 在chunk末尾附近寻找最近的句号作为切分点
            last_period = chunk.rfind("。")
            last_newline = chunk.rfind("\n")
            # 取更靠后的自然边界
            cut_point = max(last_period, last_newline)
            if cut_point > chunk_size * 0.5:  # 如果切分点不太靠前，则在此切分
                chunk = chunk[:cut_point + 1]
                end = start + len(chunk)

        chunks.append(chunk.strip())
        start = end - overlap  # 下一块的起始位置，向后重叠

    return chunks
