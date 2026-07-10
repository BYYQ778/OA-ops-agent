"""
文档解析 v2 — 多模态解析编排器
-------------------------------
调度 MinerU（可选）与传统解析器，按文件类型自动选择最优解析方案。

支持格式: PDF / DOCX / PPTX / HTML / TXT / 图片 (jpg/png/bmp/tiff/webp)
解析结果: ParsedDocument (text + tables + formulas + metadata)
"""

import os
import logging
from dataclasses import dataclass, field
from typing import List, Dict, Any, Optional

from utils.config import config

logger = logging.getLogger("doc_parser_v2")


@dataclass
class ParsedDocument:
    """解析后的文档结构化内容。"""
    text: str
    tables: List[str] = field(default_factory=list)        # Markdown 格式表格
    formulas: List[str] = field(default_factory=list)      # LaTeX 格式公式
    metadata: Dict[str, Any] = field(default_factory=dict)
    file_path: str = ""

    def to_embedding_text(self) -> str:
        """
        合并所有内容为适合向量化的单一文本。
        表格和公式用标记包裹，保留结构信息。
        """
        parts = [self.text]
        for i, table in enumerate(self.tables):
            parts.append(f"\n[TABLE_{i + 1}]\n{table}\n[/TABLE_{i + 1}]")
        for i, formula in enumerate(self.formulas):
            parts.append(f"\n[FORMULA_{i + 1}]\n{formula}\n[/FORMULA_{i + 1}]")
        return "\n".join(parts)


# ---- 扫描 PDF 检测 ----

def _is_scanned_pdf(file_path: str) -> bool:
    """判断 PDF 是否为扫描件（文本量极少）。"""
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        total_chars = 0
        pages_checked = 0
        for page in reader.pages:
            text = page.extract_text()
            if text:
                total_chars += len(text.strip())
            pages_checked += 1
            if pages_checked >= 3:  # 只检查前 3 页
                break
        return total_chars < 50
    except Exception:
        return False


# ---- 各格式解析器 ----

def _parse_with_mineru(file_path: str) -> Optional[ParsedDocument]:
    """使用 MinerU 解析，不可用时返回 None。"""
    from utils.mineru_adapter import try_mineru_parse
    result = try_mineru_parse(file_path)
    if result is None:
        return None
    return ParsedDocument(
        text=result.get("text", ""),
        tables=result.get("tables", []),
        formulas=result.get("formulas", []),
        metadata=result.get("metadata", {}),
        file_path=file_path,
    )


def _parse_pptx_native(file_path: str) -> str:
    """使用 python-pptx 提取 PowerPoint 文本。"""
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = para.text.strip()
                    if txt:
                        slide_text.append(txt)
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(" | ".join(cells))
                if rows:
                    slide_text.append("\n".join(rows))
        if slide_text:
            parts.append(f"[Slide {i}]\n" + "\n".join(slide_text))
    return "\n\n".join(parts)


def _parse_html_native(file_path: str) -> str:
    """使用 BeautifulSoup 提取 HTML 纯文本。"""
    from bs4 import BeautifulSoup
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        soup = BeautifulSoup(f.read(), "lxml")
    # 移除 script 和 style
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


# ---- 主解析入口 ----

def parse_document_v2(file_path: str, use_mineru: bool = None) -> str:
    """
    多模态文档解析入口。

    根据文件类型和 MinerU 可用性自动选择解析器，
    返回适合向量化和分块的纯文本。

    Args:
        file_path: 文档路径
        use_mineru: 是否尝试 MinerU。None=按 config 决定。

    Returns:
        清洗后的纯文本（含表格/公式标记）
    """
    if use_mineru is None:
        use_mineru = config.get("knowledge_base.mineru.enabled", False)

    ext = os.path.splitext(file_path)[1].lower()

    # ---- MinerU 路径 ----
    if use_mineru:
        doc = _parse_with_mineru(file_path)
        if doc is not None:
            return doc.to_embedding_text()
        logger.info("MinerU 不可用，回退传统解析器")

    # ---- 传统解析路径 ----
    from utils.doc_parser import (
        parse_pdf, parse_docx, parse_txt, parse_image, clean_text
    )

    if ext == ".pdf":
        if _is_scanned_pdf(file_path):
            # 扫描件：尝试 MinerU → OCR 回退
            if use_mineru:
                doc = _parse_with_mineru(file_path)
                if doc is not None:
                    return doc.to_embedding_text()
            logger.info(f"扫描件 PDF，使用 CnOCR: {os.path.basename(file_path)}")
            return parse_image(file_path)
        return clean_text(parse_pdf(file_path))

    elif ext == ".docx":
        return clean_text(parse_docx(file_path))

    elif ext == ".pptx":
        return clean_text(_parse_pptx_native(file_path))

    elif ext in (".htm", ".html"):
        return clean_text(_parse_html_native(file_path))

    elif ext == ".txt":
        return clean_text(parse_txt(file_path))

    elif ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp"):
        return parse_image(file_path)

    else:
        raise ValueError(
            f"不支持的文件格式: {ext}"
            f"（支持 .pdf / .docx / .pptx / .html / .txt / 图片格式）"
        )
